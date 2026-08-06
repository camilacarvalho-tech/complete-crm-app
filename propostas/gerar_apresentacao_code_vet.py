#!/usr/bin/env python3
"""
Apresentação Comercial PPTX — CODE VET
CRM + SaaS para Clínicas Veterinárias | CODE Tecnologia Empresarial

Cores oficiais CODE (site codetechoficial.com.br):
  - Fundo dark: #0A0E1A
  - Accent teal: #14B8A6
  - Amber: #F59E0B
  - Blue: #3B82F6
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand palette (CODE) ──────────────────────────────────────────
BG = RGBColor(0x0A, 0x0E, 0x1A)
BG_CARD = RGBColor(0x11, 0x16, 0x27)
BG_CARD_ALT = RGBColor(0x15, 0x1B, 0x2E)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_200 = RGBColor(0xE5, 0xE7, 0xEB)
SLATE_300 = RGBColor(0xD1, 0xD5, 0xDB)
SLATE_400 = RGBColor(0x9C, 0xA3, 0xAF)
SLATE_500 = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED_SOFT = RGBColor(0xF8, 0x71, 0x71)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 16

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "CODE_VET_Apresentacao_CRM_SaaS_Clinicas_Veterinarias.pptx"
ARTIFACT = Path("/opt/cursor/artifacts") / OUT_FILE.name
ARTIFACT_SHORT = Path("/opt/cursor/artifacts") / "CODE_VET_Apresentacao.pptx"


def _set_run_font(run, size_pt, color, bold=False, name="Calibri"):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def _fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    return shape


def _add_round_rect(slide, left, top, width, height, color: RGBColor, radius_pct=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    try:
        shape.adjustments[0] = radius_pct
    except Exception:
        pass
    return shape


def _textbox(slide, left, top, width, height, text, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, color, bold=bold, name=font)
    return box


def _multiline(slide, left, top, width, height, lines, size=14, color=SLATE_300,
               bold=False, spacing=6, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text = line[0]
            c = line[1] if len(line) > 1 else color
            b = line[2] if len(line) > 2 else bold
            sz = line[3] if len(line) > 3 else size
        else:
            text, c, b, sz = line, color, bold, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = text
        _set_run_font(run, sz, c, bold=b)
    return box


def _accent_bar(slide, left, top, width=Inches(0.1), height=Inches(0.42), color=TEAL):
    return _add_rect(slide, left, top, width, height, color)


def _section_title(slide, title, subtitle=None, y=Inches(0.32)):
    _accent_bar(slide, Inches(0.55), y + Inches(0.08))
    _textbox(slide, Inches(0.85), y, Inches(11.8), Inches(0.5), title, size=26, color=WHITE, bold=True)
    if subtitle:
        _textbox(slide, Inches(0.85), y + Inches(0.48), Inches(11.8), Inches(0.32),
                 subtitle, size=12, color=SLATE_400)


def _footer(slide, page: int):
    _add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), BG_CARD)
    _textbox(slide, Inches(0.5), Inches(7.18), Inches(9), Inches(0.28),
             "CODE VET  ·  CODE Tecnologia Empresarial  ·  codetechoficial.com.br",
             size=10, color=SLATE_500)
    _textbox(slide, Inches(10.5), Inches(7.18), Inches(2.3), Inches(0.28),
             f"{page} / {TOTAL}", size=10, color=SLATE_500, align=PP_ALIGN.RIGHT)


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    return slide


def _pill(slide, left, top, width, height, text, bg=TEAL, fg=BG, size=11):
    shape = _add_round_rect(slide, left, top, width, height, bg, radius_pct=0.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_run_font(run, size, fg, bold=True)
    return shape


def _chip(slide, left, top, width, height, text, accent=TEAL):
    _add_round_rect(slide, left, top, width, height, BG_CARD, 0.12)
    _add_rect(slide, left, top, Inches(0.07), height, accent)
    _textbox(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), height - Inches(0.15),
             text, size=12, color=SLATE_200, bold=False)


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════

def slide_01_capa(prs):
    s = _new_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), TEAL)
    _add_rect(s, 0, 0, Inches(0.12), SLIDE_H, TEAL)

    _textbox(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.35),
             "CODE Tecnologia Empresarial", size=14, color=SLATE_400)
    _pill(s, Inches(0.9), Inches(1.4), Inches(2.4), Inches(0.38),
          "CODE VET  ·  SaaS", bg=TEAL, fg=BG, size=12)

    _textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.7),
             "CODE VET", size=48, color=WHITE, bold=True)
    _textbox(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.5),
             "Plataforma Inteligente para Clínicas Veterinárias", size=22, color=TEAL)
    _textbox(s, Inches(0.9), Inches(3.5), Inches(11), Inches(0.4),
             "CRM  ·  ERP  ·  Inteligência Artificial  ·  Automação  ·  Business Intelligence",
             size=14, color=SLATE_300)

    _add_rect(s, 0, Inches(5.35), SLIDE_W, Inches(2.15), BG_CARD)
    _textbox(s, Inches(0.9), Inches(5.6), Inches(7), Inches(0.3),
             "Apresentado por", size=11, color=SLATE_500)
    _textbox(s, Inches(0.9), Inches(5.9), Inches(7), Inches(0.4),
             "Camila Carvalho", size=20, color=WHITE, bold=True)
    _textbox(s, Inches(0.9), Inches(6.35), Inches(7), Inches(0.3),
             "Bacharel em Ciência de Dados  ·  CODE Tecnologia Empresarial",
             size=12, color=SLATE_400)
    _textbox(s, Inches(0.9), Inches(6.75), Inches(8), Inches(0.3),
             "codetechoficial.com.br  ·  (14) 99690-2902", size=12, color=TEAL)

    _textbox(s, Inches(9.2), Inches(6.0), Inches(3.5), Inches(0.8),
             "Apresentação Comercial\nProduto SaaS Veterinária",
             size=12, color=SLATE_400, align=PP_ALIGN.RIGHT)


def slide_02_quem_somos(prs):
    s = _new_slide(prs)
    _section_title(s, "Quem é a CODE",
                   "Plataformas SaaS modernas para automatizar, reduzir custos e aumentar produtividade")
    _footer(s, 2)

    # Left visual panel (programming / tech feel)
    _add_round_rect(s, Inches(0.6), Inches(1.35), Inches(4.5), Inches(5.4), BG_CARD, 0.05)
    _add_rect(s, Inches(0.6), Inches(1.35), Inches(4.5), Inches(0.1), TEAL)
    _textbox(s, Inches(0.9), Inches(1.7), Inches(3.9), Inches(0.4),
             "</>  Engenharia de Software", size=14, color=TEAL, bold=True)
    code_lines = [
        ("const code = {", TEAL, False, 12),
        ("  crm: true,", SLATE_300, False, 12),
        ("  erp: true,", SLATE_300, False, 12),
        ("  ai: 'CODE AI',", AMBER, False, 12),
        ("  cloud: '24/7',", BLUE, False, 12),
        ("  focus: 'resultado'", GREEN, False, 12),
        ("}", TEAL, False, 12),
        ("", SLATE_500, False, 8),
        ("// +300 empresas", SLATE_500, False, 11),
        ("// SaaS · Dados · IA", SLATE_500, False, 11),
    ]
    _multiline(s, Inches(0.95), Inches(2.3), Inches(3.8), Inches(4.0), code_lines, spacing=4)

    _textbox(s, Inches(5.4), Inches(1.4), Inches(7.3), Inches(1.2),
             "A CODE Tecnologia Empresarial desenvolve plataformas SaaS modernas "
             "para empresas que desejam automatizar processos, reduzir custos e "
             "aumentar a produtividade.",
             size=14, color=SLATE_300)

    specs = [
        ("CRM", TEAL), ("ERP", BLUE), ("Inteligência Artificial", PURPLE),
        ("Ciência de Dados", AMBER), ("Desenvolvimento Web", TEAL),
        ("Aplicativos", BLUE), ("Dashboards", GREEN), ("Cloud", PURPLE),
    ]
    for i, (label, accent) in enumerate(specs):
        col, row = i % 2, i // 2
        x = Inches(5.4) + col * Inches(3.7)
        y = Inches(2.85) + row * Inches(0.85)
        _chip(s, x, y, Inches(3.5), Inches(0.7), label, accent=accent)


def slide_03_desafio(prs):
    s = _new_slide(prs)
    _section_title(s, "O desafio das clínicas veterinárias",
                   "Operação intensa, canais fragmentados e perda de receita por falta de integração")
    _footer(s, 3)

    pains = [
        ("Agenda manual", RED_SOFT),
        ("Muitos atendimentos pelo WhatsApp", AMBER),
        ("Prontuários em papel", RED_SOFT),
        ("Controle financeiro descentralizado", AMBER),
        ("Controle de estoque difícil", AMBER),
        ("Vacinas sem acompanhamento", RED_SOFT),
        ("Clientes esquecem consultas", AMBER),
        ("Falta de indicadores", RED_SOFT),
        ("Processos repetitivos", AMBER),
        ("Perda de receita", RED_SOFT),
    ]
    for i, (label, accent) in enumerate(pains):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.35) + row * Inches(1.0)
        _add_round_rect(s, x, y, Inches(6.05), Inches(0.85), BG_CARD, 0.1)
        _add_rect(s, x, y, Inches(0.09), Inches(0.85), accent)
        _textbox(s, x + Inches(0.35), y + Inches(0.22), Inches(5.4), Inches(0.4),
                 f"⚠  {label}", size=14, color=SLATE_200)


def slide_04_solucao(prs):
    s = _new_slide(prs)
    _section_title(s, "Nossa solução — CODE VET SaaS",
                   "Uma plataforma completa para administrar toda a clínica em um único lugar")
    _footer(s, 4)

    _textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
             "Integração entre todos os setores. Um sistema. Uma operação. Uma visão.",
             size=15, color=SLATE_300)

    modules = [
        ("Atendimento", TEAL), ("Financeiro", BLUE), ("Estoque", AMBER), ("CRM", PURPLE),
        ("Agenda", TEAL), ("IA", BLUE), ("Dashboards", GREEN), ("Aplicativo", PURPLE),
    ]
    for i, (label, accent) in enumerate(modules):
        col, row = i % 4, i // 4
        x = Inches(0.7) + col * Inches(3.1)
        y = Inches(2.1) + row * Inches(2.1)
        _add_round_rect(s, x, y, Inches(2.9), Inches(1.85), BG_CARD, 0.08)
        _add_rect(s, x, y, Inches(2.9), Inches(0.1), accent)
        _textbox(s, x + Inches(0.2), y + Inches(0.55), Inches(2.5), Inches(0.4),
                 "✔", size=22, color=accent, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.15), y + Inches(1.05), Inches(2.6), Inches(0.45),
                 label, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_05_dashboard(prs):
    s = _new_slide(prs)
    _section_title(s, "Dashboard Inteligente",
                   "Indicadores em tempo real para decisões rápidas e gestão baseada em dados")
    _footer(s, 5)

    kpis = [
        ("Receita", "R$", TEAL), ("Lucro", "%", GREEN), ("Consultas", "#", BLUE),
        ("Cirurgias", "#", PURPLE), ("Vacinas", "#", AMBER), ("Internações", "#", TEAL),
        ("Novos Clientes", "+", BLUE), ("Clientes Ativos", "#", GREEN),
        ("Ticket Médio", "R$", AMBER), ("Fluxo de Caixa", "R$", TEAL),
        ("Agendamentos", "#", BLUE), ("Cancelamentos", "!", RED_SOFT),
    ]
    for i, (label, unit, accent) in enumerate(kpis):
        col, row = i % 4, i // 4
        x = Inches(0.6) + col * Inches(3.15)
        y = Inches(1.35) + row * Inches(1.75)
        _add_round_rect(s, x, y, Inches(3.0), Inches(1.55), BG_CARD, 0.08)
        _add_rect(s, x, y, Inches(3.0), Inches(0.08), accent)
        _textbox(s, x + Inches(0.2), y + Inches(0.3), Inches(2.6), Inches(0.35),
                 unit, size=12, color=accent, bold=True)
        _textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(2.6), Inches(0.5),
                 label, size=15, color=WHITE, bold=True)


def slide_06_crm(prs):
    s = _new_slide(prs)
    _section_title(s, "CRM Veterinário — o coração do sistema",
                   "Cadastro completo de tutores, pets e histórico clínico-financeiro")
    _footer(s, 6)

    columns = [
        ("Tutor", TEAL, [
            "Nome", "CPF", "Telefone", "WhatsApp", "E-mail",
            "Endereço", "Histórico Financeiro",
        ]),
        ("Pet", BLUE, [
            "Nome", "Espécie / Raça", "Sexo / Peso / Idade", "Cor / Foto",
            "Microchip", "Castração", "Alergias / Medicamentos", "Vacinas",
        ]),
        ("Histórico Médico", AMBER, [
            "Consultas", "Exames", "Vacinas", "Internações",
            "Cirurgias", "Receitas", "Fotos / Arquivos", "Anexos",
        ]),
    ]
    card_w = Inches(3.95)
    for i, (title, accent, items) in enumerate(columns):
        x = Inches(0.6) + i * (card_w + Inches(0.2))
        _add_round_rect(s, x, Inches(1.3), card_w, Inches(5.5), BG_CARD, 0.05)
        _add_rect(s, x, Inches(1.3), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.25), Inches(1.6), card_w - Inches(0.4), Inches(0.4),
                 title, size=16, color=WHITE, bold=True)
        _multiline(s, x + Inches(0.25), Inches(2.2), card_w - Inches(0.4), Inches(4.3),
                   [f"•  {it}" for it in items], size=13, color=SLATE_300, spacing=8)


def slide_07_agenda(prs):
    s = _new_slide(prs)
    _section_title(s, "Agenda Inteligente",
                   "Organização por profissional, sala e procedimento — com automações")
    _footer(s, 7)

    items = [
        ("Agenda por médico veterinário", TEAL),
        ("Agenda por sala", BLUE),
        ("Agenda por procedimento", PURPLE),
        ("Confirmação automática", GREEN),
        ("Lembretes inteligentes", AMBER),
        ("Reagendamento facilitado", TEAL),
        ("Fila de espera", BLUE),
        ("Integração Google Calendar", PURPLE),
    ]
    for i, (label, accent) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.4) + row * Inches(1.25)
        _add_round_rect(s, x, y, Inches(6.05), Inches(1.1), BG_CARD, 0.1)
        _add_rect(s, x, y, Inches(0.1), Inches(1.1), accent)
        _textbox(s, x + Inches(0.4), y + Inches(0.32), Inches(5.4), Inches(0.45),
                 f"📅  {label}", size=16, color=WHITE, bold=True)


def slide_08_prontuario(prs):
    s = _new_slide(prs)
    _section_title(s, "Prontuário Digital",
                   "Registro clínico completo, seguro e acessível em qualquer momento")
    _footer(s, 8)

    items = [
        ("Consultas", "Anamnese, evolução e condutas"),
        ("Diagnóstico", "CID Veterinário padronizado"),
        ("Prescrição", "Receituário digital estruturado"),
        ("Exames", "Laboratório e laudos anexados"),
        ("Imagens", "Raio-X, ultrassom e arquivos"),
        ("Histórico", "Linha do tempo do paciente"),
        ("Receituário", "Emissão e controle de receitas"),
        ("Assinatura Digital", "Validade e rastreabilidade"),
    ]
    for i, (title, desc) in enumerate(items):
        col, row = i % 4, i // 4
        x = Inches(0.55) + col * Inches(3.15)
        y = Inches(1.4) + row * Inches(2.5)
        _add_round_rect(s, x, y, Inches(3.0), Inches(2.25), BG_CARD, 0.08)
        _add_rect(s, x, y, Inches(3.0), Inches(0.1), TEAL if row == 0 else BLUE)
        _textbox(s, x + Inches(0.2), y + Inches(0.5), Inches(2.6), Inches(0.5),
                 title, size=15, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.2), y + Inches(1.15), Inches(2.6), Inches(0.7),
                 desc, size=12, color=SLATE_400)


def slide_09_financeiro(prs):
    s = _new_slide(prs)
    _section_title(s, "Financeiro",
                   "Controle completo do caixa, faturamento, comissões e indicadores")
    _footer(s, 9)

    items = [
        "Contas a pagar", "Contas a receber", "Fluxo de caixa", "Mensalidades",
        "Convênios", "PIX", "Cartão", "Boletos",
        "NF-e", "DRE", "Comissões", "Relatórios",
    ]
    accents = [TEAL, BLUE, AMBER, PURPLE] * 3
    for i, label in enumerate(items):
        col, row = i % 4, i // 4
        x = Inches(0.6) + col * Inches(3.15)
        y = Inches(1.4) + row * Inches(1.7)
        _add_round_rect(s, x, y, Inches(3.0), Inches(1.5), BG_CARD, 0.1)
        _add_rect(s, x, y, Inches(0.09), Inches(1.5), accents[i])
        _textbox(s, x + Inches(0.3), y + Inches(0.5), Inches(2.5), Inches(0.5),
                 f"💰  {label}", size=14, color=WHITE, bold=True)


def slide_10_estoque(prs):
    s = _new_slide(prs)
    _section_title(s, "Estoque Inteligente",
                   "Controle de insumos críticos com validade, lote e reposição automática")
    _footer(s, 10)

    left = [
        ("Medicamentos", TEAL),
        ("Vacinas", BLUE),
        ("Produtos", PURPLE),
        ("Ração", AMBER),
        ("Materiais", GREEN),
    ]
    right = [
        ("Controle por lote", TEAL),
        ("Validade", RED_SOFT),
        ("Fornecedor", BLUE),
        ("Compra automática", AMBER),
        ("Estoque mínimo", PURPLE),
        ("Código de barras", GREEN),
    ]

    _textbox(s, Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.35),
             "Categorias", size=14, color=TEAL, bold=True)
    for i, (label, accent) in enumerate(left):
        y = Inches(1.8) + i * Inches(0.95)
        _add_round_rect(s, Inches(0.7), y, Inches(5.8), Inches(0.85), BG_CARD, 0.1)
        _add_rect(s, Inches(0.7), y, Inches(0.09), Inches(0.85), accent)
        _textbox(s, Inches(1.05), y + Inches(0.22), Inches(5.2), Inches(0.4),
                 label, size=15, color=WHITE, bold=True)

    _textbox(s, Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.35),
             "Controles", size=14, color=AMBER, bold=True)
    for i, (label, accent) in enumerate(right):
        y = Inches(1.8) + i * Inches(0.8)
        _add_round_rect(s, Inches(6.9), y, Inches(5.8), Inches(0.7), BG_CARD, 0.12)
        _add_rect(s, Inches(6.9), y, Inches(0.09), Inches(0.7), accent)
        _textbox(s, Inches(7.25), y + Inches(0.15), Inches(5.2), Inches(0.4),
                 f"✓  {label}", size=14, color=SLATE_200)


def slide_11_ia(prs):
    s = _new_slide(prs)
    _section_title(s, "Inteligência Artificial — CODE AI",
                   "Automação inteligente para atendimento, operação e crescimento")
    _footer(s, 11)

    abilities = [
        ("Responder clientes", TEAL),
        ("Triagem inicial", BLUE),
        ("Agendar consultas", PURPLE),
        ("Responder WhatsApp", GREEN),
        ("Analisar indicadores", AMBER),
        ("Criar campanhas", TEAL),
        ("Enviar lembretes", BLUE),
        ("Prever demanda", PURPLE),
        ("Gerar relatórios", AMBER),
        ("Analisar iniciativas", GREEN),
    ]
    for i, (label, accent) in enumerate(abilities):
        col, row = i % 5, i // 5
        x = Inches(0.55) + col * Inches(2.5)
        y = Inches(1.5) + row * Inches(2.5)
        _add_round_rect(s, x, y, Inches(2.35), Inches(2.2), BG_CARD, 0.08)
        _add_rect(s, x, y, Inches(2.35), Inches(0.1), accent)
        _textbox(s, x + Inches(0.15), y + Inches(0.7), Inches(2.05), Inches(0.4),
                 "AI", size=18, color=accent, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.15), y + Inches(1.2), Inches(2.05), Inches(0.7),
                 label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_12_automacoes(prs):
    s = _new_slide(prs)
    _section_title(s, "Automações",
                   "Comunicação e marketing contínuos sem retrabalho da equipe")
    _footer(s, 12)

    channels = [
        ("Canais", TEAL, ["WhatsApp", "SMS", "E-mail"]),
        ("Lembretes", BLUE, ["Vacina", "Retorno", "Aniversário do Pet"]),
        ("Relacionamento", PURPLE, ["Pesquisa de satisfação", "Cobranças", "Marketing"]),
        ("Crescimento", AMBER, ["Campanhas", "Remarketing", "Reativação"]),
    ]
    card_w = Inches(2.95)
    for i, (title, accent, items) in enumerate(channels):
        x = Inches(0.55) + i * (card_w + Inches(0.2))
        _add_round_rect(s, x, Inches(1.4), card_w, Inches(5.3), BG_CARD, 0.06)
        _add_rect(s, x, Inches(1.4), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.2), Inches(1.8), card_w - Inches(0.35), Inches(0.45),
                 title, size=16, color=WHITE, bold=True)
        for j, item in enumerate(items):
            y = Inches(2.6) + j * Inches(1.1)
            _add_round_rect(s, x + Inches(0.2), y, card_w - Inches(0.4), Inches(0.9), BG_CARD_ALT, 0.15)
            _textbox(s, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.7), Inches(0.4),
                     item, size=13, color=SLATE_200, align=PP_ALIGN.CENTER)


def slide_13_app(prs):
    s = _new_slide(prs)
    _section_title(s, "Aplicativo Mobile",
                   "Experiência digital para a clínica e para o tutor — codetechoficial.com.br")
    _footer(s, 13)

    # Phone-like visual panel
    _add_round_rect(s, Inches(0.8), Inches(1.4), Inches(4.2), Inches(5.3), BG_CARD, 0.08)
    _add_rect(s, Inches(0.8), Inches(1.4), Inches(4.2), Inches(0.1), TEAL)
    _textbox(s, Inches(1.1), Inches(1.8), Inches(3.6), Inches(0.4),
             "App da Clínica", size=18, color=WHITE, bold=True)
    _textbox(s, Inches(1.1), Inches(2.3), Inches(3.6), Inches(0.4),
             "Tutor acompanha tudo no celular", size=12, color=SLATE_400)
    app_items = ["Vacinas", "Exames", "Receitas", "Agendamento",
                 "Carteira Digital do Pet", "Notificações", "Chat"]
    _multiline(s, Inches(1.1), Inches(2.9), Inches(3.6), Inches(3.5),
               [f"✓  {it}" for it in app_items], size=14, color=SLATE_200, spacing=10)

    rights = [
        ("Carteira Digital do Pet", "Histórico vacinal e documentos sempre à mão"),
        ("Agendamento online", "Marque e remarque com poucos toques"),
        ("Chat & Notificações", "Comunicação direta com a clínica"),
        ("Exames e Receitas", "Acesso rápido a resultados e prescrições"),
    ]
    for i, (title, desc) in enumerate(rights):
        y = Inches(1.4) + i * Inches(1.3)
        _add_round_rect(s, Inches(5.4), y, Inches(7.2), Inches(1.15), BG_CARD, 0.1)
        _add_rect(s, Inches(5.4), y, Inches(0.1), Inches(1.15), TEAL if i % 2 == 0 else BLUE)
        _textbox(s, Inches(5.8), y + Inches(0.2), Inches(6.5), Inches(0.35),
                 title, size=15, color=WHITE, bold=True)
        _textbox(s, Inches(5.8), y + Inches(0.6), Inches(6.5), Inches(0.35),
                 desc, size=12, color=SLATE_400)


def slide_14_beneficios(prs):
    s = _new_slide(prs)
    _section_title(s, "Benefícios",
                   "Resultados práticos na operação, no atendimento e no faturamento")
    _footer(s, 14)

    benefits = [
        ("Redução de custos", TEAL),
        ("Mais produtividade", BLUE),
        ("Mais organização", PURPLE),
        ("Menos retrabalho", AMBER),
        ("Atendimento rápido", GREEN),
        ("Mais operação", TEAL),
        ("Clientes fidelizados", BLUE),
        ("Gestão completa", PURPLE),
        ("Segurança dos dados", AMBER),
        ("Escalabilidade", GREEN),
    ]
    for i, (label, accent) in enumerate(benefits):
        col, row = i % 5, i // 5
        x = Inches(0.55) + col * Inches(2.5)
        y = Inches(1.6) + row * Inches(2.5)
        _add_round_rect(s, x, y, Inches(2.35), Inches(2.2), BG_CARD, 0.08)
        _add_rect(s, x, y, Inches(2.35), Inches(0.1), accent)
        _textbox(s, x + Inches(0.15), y + Inches(0.7), Inches(2.05), Inches(0.4),
                 "★", size=20, color=accent, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.15), y + Inches(1.2), Inches(2.05), Inches(0.7),
                 label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_15_roteiro(prs):
    s = _new_slide(prs)
    _section_title(s, "Roteiro do Projeto",
                   "Do levantamento de requisitos à implantação e suporte contínuo")
    _footer(s, 15)

    phases = [
        ("Fase 1", "Descoberta", "Levantamento de requisitos\nProtótipos (UI/UX)\nAprovação do fluxo", TEAL),
        ("Fase 2", "Core CRM", "Desenvolvimento do CRM\nAgenda\nTutores e pets\nProntuário digital", BLUE),
        ("Fase 3", "Gestão", "Financeiro\nEstoque\nRelatórios\nDashboard", PURPLE),
        ("Fase 4", "Inteligência", "CODE AI\nWhatsApp\nAutomações\nApp Mobile", AMBER),
        ("Fase 5", "Go-Live", "Testes\nImplantação\nTreinamento\nSuporte contínuo", GREEN),
    ]
    card_w = Inches(2.35)
    for i, (phase, title, body, accent) in enumerate(phases):
        x = Inches(0.5) + i * (card_w + Inches(0.15))
        _add_round_rect(s, x, Inches(1.35), card_w, Inches(5.4), BG_CARD, 0.06)
        _add_rect(s, x, Inches(1.35), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.15), Inches(1.65), card_w - Inches(0.25), Inches(0.35),
                 phase, size=12, color=accent, bold=True)
        _textbox(s, x + Inches(0.15), Inches(2.05), card_w - Inches(0.25), Inches(0.4),
                 title, size=15, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.15), Inches(2.7), card_w - Inches(0.25), Inches(3.6),
                 body, size=12, color=SLATE_300)


def slide_16_cta(prs):
    s = _new_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), TEAL)
    _add_rect(s, 0, 0, Inches(0.12), SLIDE_H, TEAL)
    _footer(s, 16)

    _pill(s, Inches(4.55), Inches(1.2), Inches(4.2), Inches(0.4),
          "VAMOS TRANSFORMAR SUA CLÍNICA?", bg=TEAL, fg=BG, size=12)

    _textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.6),
             "CODE Tecnologia Empresarial", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.45),
             "Tecnologia que impulsiona empresas.", size=16, color=TEAL, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.4),
             "CRM  ·  ERP  ·  IA  ·  Automação  ·  App  ·  Business Intelligence",
             size=13, color=SLATE_400, align=PP_ALIGN.CENTER)

    contacts = [
        ("Site", "codetechoficial.com.br"),
        ("WhatsApp", "(14) 99690-2902"),
        ("Apresentação", "Camila Carvalho"),
    ]
    for i, (label, value) in enumerate(contacts):
        x = Inches(1.4) + i * Inches(3.85)
        _add_round_rect(s, x, Inches(3.9), Inches(3.5), Inches(1.5), BG_CARD, 0.1)
        _textbox(s, x + Inches(0.2), Inches(4.15), Inches(3.1), Inches(0.35),
                 label, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.2), Inches(4.6), Inches(3.1), Inches(0.5),
                 value, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _textbox(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.35),
             "Bacharel em Ciência de Dados  ·  CODE Tecnologia Empresarial",
             size=12, color=SLATE_500, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.5),
             "Não entregamos apenas um sistema — entregamos uma plataforma SaaS completa\n"
             "para gestão veterinária: CRM, ERP, prontuário, agenda, financeiro, estoque, IA e app.",
             size=12, color=SLATE_400, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_capa(prs)
    slide_02_quem_somos(prs)
    slide_03_desafio(prs)
    slide_04_solucao(prs)
    slide_05_dashboard(prs)
    slide_06_crm(prs)
    slide_07_agenda(prs)
    slide_08_prontuario(prs)
    slide_09_financeiro(prs)
    slide_10_estoque(prs)
    slide_11_ia(prs)
    slide_12_automacoes(prs)
    slide_13_app(prs)
    slide_14_beneficios(prs)
    slide_15_roteiro(prs)
    slide_16_cta(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    ARTIFACT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ARTIFACT))
    prs.save(str(ARTIFACT_SHORT))
    print(f"Saved: {OUT_FILE}")
    print(f"Artifact: {ARTIFACT}")
    print(f"Artifact short: {ARTIFACT_SHORT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
