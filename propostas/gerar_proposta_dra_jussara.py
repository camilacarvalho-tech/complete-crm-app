#!/usr/bin/env python3
"""
Proposta Comercial PPTX — Dra. Jussara Bertoloto
Projeto de Presença Digital | CODE Tecnologia Empresarial

Cores oficiais CODE (site codetechoficial.com.br):
  - Fundo dark: #0A0E1A
  - Accent teal: #14B8A6
  - Amber: #F59E0B
  - Blue: #3B82F6
  - Slate text: #9CA3AF / #D1D5DB
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand palette (CODE) ──────────────────────────────────────────
BG = RGBColor(0x0A, 0x0E, 0x1A)
BG_CARD = RGBColor(0x11, 0x16, 0x27)
BG_CARD_ALT = RGBColor(0x15, 0x1B, 0x2E)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
TEAL_DARK = RGBColor(0x0D, 0x94, 0x88)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_200 = RGBColor(0xE5, 0xE7, 0xEB)
SLATE_300 = RGBColor(0xD1, 0xD5, 0xDB)
SLATE_400 = RGBColor(0x9C, 0xA3, 0xAF)
SLATE_500 = RGBColor(0x6B, 0x72, 0x80)
RED_SOFT = RGBColor(0xF8, 0x71, 0x71)
GREEN = RGBColor(0x34, 0xD3, 0x99)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "CODE_Proposta_Presenca_Digital_Dra_Jussara_Bertoloto.pptx"
ARTIFACT = Path("/opt/cursor/artifacts") / OUT_FILE.name


def _set_run_font(run, size_pt, color, bold=False, name="Calibri"):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def _fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    return shape


def _add_round_rect(slide, left, top, width, height, color: RGBColor, radius_pct=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, color)
    # Soften corners a bit via adj
    try:
        shape.adjustments[0] = radius_pct
    except Exception:
        pass
    return shape


def _textbox(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, color, bold=bold, name=font)
    return box


def _multiline(slide, left, top, width, height, lines, size=14, color=SLATE_300, bold=False, spacing=6, align=PP_ALIGN.LEFT):
    """lines: list of str or (str, color, bold) tuples"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, c, b = line[0], line[1] if len(line) > 1 else color, line[2] if len(line) > 2 else bold
        else:
            text, c, b = line, color, bold
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = text
        _set_run_font(run, size, c, bold=b)
    return box


def _accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.45), color=TEAL):
    return _add_rect(slide, left, top, width, height, color)


def _section_title(slide, title, subtitle=None, y=Inches(0.35)):
    _accent_bar(slide, Inches(0.6), y + Inches(0.08), Inches(0.1), Inches(0.42))
    _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.55), title, size=28, color=WHITE, bold=True)
    if subtitle:
        _textbox(slide, Inches(0.9), y + Inches(0.5), Inches(11.5), Inches(0.35), subtitle, size=13, color=SLATE_400)


def _footer(slide, page: int, total: int = 12):
    _add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), BG_CARD)
    _textbox(slide, Inches(0.5), Inches(7.18), Inches(8), Inches(0.28),
             "CODE Tecnologia Empresarial  ·  codetechoficial.com.br", size=10, color=SLATE_500)
    _textbox(slide, Inches(10.5), Inches(7.18), Inches(2.3), Inches(0.28),
             f"{page} / {total}", size=10, color=SLATE_500, align=PP_ALIGN.RIGHT)


def _new_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    return slide


def _pill(slide, left, top, width, height, text, bg=TEAL, fg=BG, size=11):
    shape = _add_round_rect(slide, left, top, width, height, bg, radius_pct=0.5)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_run_font(run, size, fg, bold=True)
    try:
        tf.paragraphs[0].space_before = Pt(4)
    except Exception:
        pass
    return shape


def _card(slide, left, top, width, height, title, body_lines, accent=TEAL, icon=None):
    _add_round_rect(slide, left, top, width, height, BG_CARD, radius_pct=0.06)
    # left accent strip
    _add_rect(slide, left, top, Inches(0.08), height, accent)
    title_text = f"{icon}  {title}" if icon else title
    _textbox(slide, left + Inches(0.25), top + Inches(0.18), width - Inches(0.4), Inches(0.4),
             title_text, size=14, color=WHITE, bold=True)
    if body_lines:
        _multiline(slide, left + Inches(0.25), top + Inches(0.6), width - Inches(0.4), height - Inches(0.75),
                   body_lines, size=12, color=SLATE_300, spacing=4)


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════

def slide_cover(prs):
    s = _new_slide(prs)
    # Top teal line
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), TEAL)
    # Side accent
    _add_rect(s, 0, 0, Inches(0.12), SLIDE_H, TEAL)

    _pill(s, Inches(0.9), Inches(1.4), Inches(3.2), Inches(0.38),
          "PROJETO DE PRESENÇA DIGITAL", bg=TEAL, fg=BG, size=11)

    _textbox(s, Inches(0.9), Inches(2.1), Inches(11), Inches(0.9),
             "Dra. Jussara Bertoloto", size=42, color=WHITE, bold=True)
    _textbox(s, Inches(0.9), Inches(3.0), Inches(11), Inches(0.5),
             "Especialista em Tratamento de Feridas", size=22, color=TEAL, bold=False)
    _textbox(s, Inches(0.9), Inches(3.6), Inches(10), Inches(0.6),
             "Proposta comercial de Landing Page Premium + Estratégia Digital\npara geração de pacientes e autoridade médica online.",
             size=14, color=SLATE_400)

    # Bottom brand block
    _add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(1.9), BG_CARD)
    _textbox(s, Inches(0.9), Inches(5.85), Inches(6), Inches(0.35),
             "Desenvolvido por", size=11, color=SLATE_500)
    _textbox(s, Inches(0.9), Inches(6.15), Inches(8), Inches(0.45),
             "CODE Tecnologia Empresarial", size=22, color=WHITE, bold=True)
    _textbox(s, Inches(0.9), Inches(6.65), Inches(8), Inches(0.35),
             "codetechoficial.com.br  ·  (14) 99690-2902", size=12, color=TEAL)

    _textbox(s, Inches(9.5), Inches(6.2), Inches(3.2), Inches(0.7),
             "Apresentação Comercial\n2026", size=12, color=SLATE_400, align=PP_ALIGN.RIGHT)


def slide_quem_somos(prs):
    s = _new_slide(prs)
    _section_title(s, "Quem Somos — CODE Tecnologia Empresarial",
                   "Ecossistemas completos de tecnologia, automação e marketing de alta conversão")
    _footer(s, 3)

    intro = ("A CODE Tecnologia Empresarial é especializada em ecossistemas completos de "
             "tecnologia, automação e marketing de alta conversão.")
    _textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.55), intro, size=14, color=SLATE_300)

    services = [
        ("Desenvolvimento de Sistemas", "Nexus CRM / ERP e plataformas sob medida", BLUE, "01"),
        ("Landing Pages & Sites Premium", "Alta conversão e presença institucional", TEAL, "02"),
        ("Automação & Inteligência Artificial", "Processos inteligentes e atendimento com IA", PURPLE, "03"),
        ("Gestão de Tráfego Pago", "Google Ads & Meta Ads com foco em ROI", AMBER, "04"),
        ("SEO Estratégico", "Posicionamento no topo do Google", GREEN, "05"),
        ("Suporte Contínuo 24/7", "Acompanhamento técnico e comercial", TEAL, "06"),
    ]

    cols, rows = 3, 2
    card_w, card_h = Inches(3.85), Inches(1.55)
    gap_x, gap_y = Inches(0.25), Inches(0.22)
    start_x, start_y = Inches(0.7), Inches(2.1)

    for i, (title, desc, accent, num) in enumerate(services):
        r, c = divmod(i, cols)
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        _add_round_rect(s, x, y, card_w, card_h, BG_CARD, 0.06)
        _add_rect(s, x, y, Inches(0.08), card_h, accent)
        _textbox(s, x + Inches(0.25), y + Inches(0.25), Inches(0.6), Inches(0.35),
                 num, size=16, color=accent, bold=True)
        _textbox(s, x + Inches(0.25), y + Inches(0.6), card_w - Inches(0.4), Inches(0.35),
                 title, size=13, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.25), y + Inches(1.0), card_w - Inches(0.4), Inches(0.4),
                 desc, size=11, color=SLATE_400)

    # Stats bar
    stats = [
        ("+300", "empresas transformadas"),
        ("+50K", "atendimentos"),
        ("98%", "de satisfação"),
        ("24/7", "suporte disponível"),
    ]
    bar_y = Inches(5.55)
    _add_round_rect(s, Inches(0.7), bar_y, Inches(11.9), Inches(1.25), BG_CARD_ALT, 0.05)
    for i, (n, label) in enumerate(stats):
        x = Inches(1.0) + i * Inches(3.0)
        _textbox(s, x, bar_y + Inches(0.25), Inches(2.6), Inches(0.45), n, size=26, color=TEAL, bold=True)
        _textbox(s, x, bar_y + Inches(0.7), Inches(2.6), Inches(0.35), label, size=11, color=SLATE_400)


def slide_problema(prs):
    s = _new_slide(prs)
    _section_title(s, "O Problema de Mercado",
                   "Pacientes e familiares buscam soluções urgentes no Google — todos os dias")
    _footer(s, 4)

    searches = [
        '"Curativo especializado perto de mim"',
        '"Tratamento para pé diabético e escaras"',
        '"Enfermeira especialista em feridas a domicílio"',
        '"Avaliação de lesões de difícil cicatrização"',
    ]

    _textbox(s, Inches(0.7), Inches(1.4), Inches(6.2), Inches(0.4),
             "Buscas reais de pacientes", size=14, color=TEAL, bold=True)

    for i, q in enumerate(searches):
        y = Inches(1.9) + i * Inches(0.85)
        _add_round_rect(s, Inches(0.7), y, Inches(6.4), Inches(0.7), BG_CARD, 0.1)
        _textbox(s, Inches(0.95), y + Inches(0.18), Inches(5.9), Inches(0.4),
                 f"🔍  {q}", size=13, color=SLATE_200)

    # Consequence panel
    _add_round_rect(s, Inches(7.5), Inches(1.9), Inches(5.1), Inches(4.3), BG_CARD, 0.05)
    _add_rect(s, Inches(7.5), Inches(1.9), Inches(5.1), Inches(0.1), AMBER)
    _textbox(s, Inches(7.8), Inches(2.25), Inches(4.5), Inches(0.4),
             "A Consequência", size=18, color=AMBER, bold=True)
    _multiline(s, Inches(7.8), Inches(2.85), Inches(4.5), Inches(3.0), [
        "Sem uma presença digital otimizada e de alta autoridade, esses pacientes encontram concorrentes ou serviços genéricos.",
        "",
        "Cada dia sem posicionamento é uma oportunidade de atendimento e agendamento perdida.",
        "",
        "→ Visibilidade zero no momento de decisão",
        "→ Agenda vazia por falta de demanda captada",
        "→ Autoridade médica não comunicada online",
    ], size=13, color=SLATE_300, spacing=6)


def slide_servicos(prs):
    s = _new_slide(prs)
    _section_title(s, "Produtos e Serviços da Dra. Jussara",
                   "Mapeamento completo dos atendimentos, tratamentos e procedimentos especializados")
    _footer(s, 5)

    cats = [
        ("Tratamento de Lesões Complexas", TEAL, [
            "• Tratamento especializado para Pé Diabético",
            "• Cuidado avançado de Escaras (Úlceras por Pressão)",
            "• Tratamento de Úlceras Venosas e Arteriais",
            "• Cuidados com Queimaduras e Feridas Cirúrgicas",
        ]),
        ("Terapia Avançada & Coberturas", BLUE, [
            "• Aplicação de curativos especiais e tecnologia",
            "  de cicatrização acelerada",
            "• Desbridamento (remoção de tecido desvitalizado)",
            "• Avaliação e acompanhamento contínuo da lesão",
        ]),
        ("Modalidades de Atendimento", AMBER, [
            "• Atendimento Domiciliar (Home Care)",
            "  Conforto e segurança no lar do paciente",
            "• Consultório Especializado",
            "  Infraestrutura completa para procedimentos",
        ]),
    ]

    card_w = Inches(3.9)
    for i, (title, accent, lines) in enumerate(cats):
        x = Inches(0.7) + i * (card_w + Inches(0.25))
        _add_round_rect(s, x, Inches(1.45), card_w, Inches(5.2), BG_CARD, 0.05)
        _add_rect(s, x, Inches(1.45), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.25), Inches(1.8), card_w - Inches(0.4), Inches(0.7),
                 title, size=15, color=WHITE, bold=True)
        _multiline(s, x + Inches(0.25), Inches(2.7), card_w - Inches(0.4), Inches(3.6),
                   lines, size=12, color=SLATE_300, spacing=8)


def slide_solucao(prs):
    s = _new_slide(prs)
    _section_title(s, "A Solução CODE — Landing Page + Estratégia",
                   "Estrutura digital projetada para transmitir autoridade médica e converter em agendamentos")
    _footer(s, 6)

    items = [
        ("Design Dark / Premium", "Cores modernas, elementos visuais tecnológicos e navegação fluida.", TEAL),
        ("Foco em Conversão", "Botões estratégicos de WhatsApp em pontos-chave da página.", AMBER),
        ("Infraestrutura Completa", "Hospedagem SSD de alta velocidade, certificado SSL e proteção anti-spam.", BLUE),
        ("Rastreamento de Dados", "Métricas via Google Analytics 4, Meta Pixel e Google Tag Manager.", PURPLE),
    ]

    for i, (title, desc, accent) in enumerate(items):
        y = Inches(1.45) + i * Inches(1.2)
        _add_round_rect(s, Inches(0.7), y, Inches(11.9), Inches(1.05), BG_CARD, 0.06)
        _add_rect(s, Inches(0.7), y, Inches(0.1), Inches(1.05), accent)
        num = f"0{i+1}"
        _textbox(s, Inches(1.1), y + Inches(0.25), Inches(0.7), Inches(0.5), num, size=20, color=accent, bold=True)
        _textbox(s, Inches(2.0), y + Inches(0.18), Inches(9.8), Inches(0.35), title, size=16, color=WHITE, bold=True)
        _textbox(s, Inches(2.0), y + Inches(0.55), Inches(9.8), Inches(0.35), desc, size=13, color=SLATE_400)


def slide_estrutura_lp(prs):
    s = _new_slide(prs)
    _section_title(s, "Estrutura Detalhada da Landing Page",
                   "Arquitetura de seções otimizada para autoridade e conversão")
    _footer(s, 7)

    sections = [
        ("01", "Hero Section", "Título de alto impacto + CTA imediata para WhatsApp"),
        ("02", "Sobre a Dra. Jussara", "Trajetória, qualificações e diferenciais profissionais"),
        ("03", "Serviços & Especialidades", "Exposição clara de todos os tratamentos para feridas"),
        ("04", "Atendimento Domiciliar", "Benefícios e comodidade do atendimento em casa"),
        ("05", "Depoimentos & Casos", "Prova social com relatos reais para gerar confiança"),
        ("06", "Como Funciona", "Passo a passo: da triagem inicial ao tratamento"),
        ("07", "Perguntas Frequentes", "Respostas sobre prazos, atendimentos e planos"),
        ("08", "Localização & Contato", "Maps, telefone, e-mail e formulário inteligente"),
    ]

    for i, (num, title, desc) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * Inches(6.2)
        y = Inches(1.4) + row * Inches(1.25)
        _add_round_rect(s, x, y, Inches(5.9), Inches(1.1), BG_CARD, 0.06)
        _textbox(s, x + Inches(0.25), y + Inches(0.25), Inches(0.7), Inches(0.5),
                 num, size=20, color=TEAL, bold=True)
        _textbox(s, x + Inches(1.1), y + Inches(0.2), Inches(4.5), Inches(0.35),
                 title, size=14, color=WHITE, bold=True)
        _textbox(s, x + Inches(1.1), y + Inches(0.55), Inches(4.5), Inches(0.4),
                 desc, size=12, color=SLATE_400)


def slide_trafego(prs):
    s = _new_slide(prs)
    _section_title(s, "Estratégia de Tráfego & Visibilidade",
                   "Captura de demanda urgente + construção de autoridade + SEO local")
    _footer(s, 8)

    channels = [
        ("Google Ads", "Anúncios na Busca Ativa", BLUE, [
            "Captura de demanda urgente: exibição imediata para quem pesquisa tratamento de feridas na região.",
            "Controle total do investimento: cobrança apenas por cliques efetivos e qualificados.",
        ]),
        ("Meta Ads", "Instagram & Facebook", PURPLE, [
            "Anúncios educativos & autoridade: vídeos, carrosséis de prevenção e Stories interativos.",
            "Remarketing inteligente: reimpacta visitantes da Landing Page até a conversão.",
        ]),
        ("Google Meu Negócio", "SEO Local", TEAL, [
            "Topo do Mapa do Google: perfil 100% otimizado com fotos, horários e link de agendamento.",
            "Gestão de reputação: estratégia contínua para captação de avaliações 5 estrelas.",
        ]),
    ]

    card_w = Inches(3.9)
    for i, (title, sub, accent, lines) in enumerate(channels):
        x = Inches(0.7) + i * (card_w + Inches(0.25))
        _add_round_rect(s, x, Inches(1.45), card_w, Inches(5.2), BG_CARD, 0.05)
        _add_rect(s, x, Inches(1.45), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.25), Inches(1.8), card_w - Inches(0.4), Inches(0.4),
                 title, size=18, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.25), Inches(2.25), card_w - Inches(0.4), Inches(0.35),
                 sub, size=12, color=accent, bold=True)
        _multiline(s, x + Inches(0.25), Inches(2.9), card_w - Inches(0.4), Inches(3.4),
                   lines, size=12, color=SLATE_300, spacing=12)


def slide_infra(prs):
    s = _new_slide(prs)
    _section_title(s, "Infraestrutura Técnica CODE",
                   "Performance, integrações e segurança para uma operação digital confiável")
    _footer(s, 9)

    blocks = [
        ("Tecnologia", TEAL, [
            "Landing page 100% responsiva",
            "Hospedagem SSD ultra-rápida",
            "CDN Global para velocidade",
            "Otimização de Core Web Vitals",
        ]),
        ("Integrações", BLUE, [
            "WhatsApp Business API",
            "Formulário inteligente",
            "Automações de follow-up",
            "CRM / captura de leads",
        ]),
        ("Segurança", AMBER, [
            "Backup automático",
            "Certificado SSL (HTTPS)",
            "Proteção contra ataques",
            "Monitoramento contínuo",
        ]),
    ]

    card_w = Inches(3.9)
    for i, (title, accent, items) in enumerate(blocks):
        x = Inches(0.7) + i * (card_w + Inches(0.25))
        _add_round_rect(s, x, Inches(1.5), card_w, Inches(4.9), BG_CARD, 0.05)
        _add_rect(s, x, Inches(1.5), card_w, Inches(0.1), accent)
        _textbox(s, x + Inches(0.3), Inches(1.9), card_w - Inches(0.5), Inches(0.45),
                 title, size=18, color=WHITE, bold=True)
        for j, item in enumerate(items):
            y = Inches(2.6) + j * Inches(0.75)
            _add_round_rect(s, x + Inches(0.3), y, card_w - Inches(0.6), Inches(0.6), BG_CARD_ALT, 0.15)
            _textbox(s, x + Inches(0.5), y + Inches(0.12), card_w - Inches(1.0), Inches(0.4),
                     f"✓  {item}", size=12, color=SLATE_200)


def slide_resultados(prs):
    s = _new_slide(prs)
    _section_title(s, "Resultados Esperados (em até 90 dias)",
                   "Metas de visibilidade, agenda e retorno sobre investimento")
    _footer(s, 10)

    results = [
        ("Visibilidade Máxima", "Posicionamento forte na região de atuação", "01", TEAL),
        ("Novos Pacientes", "Fluxo constante de leads e agendamentos", "02", BLUE),
        ("Autoridade Médica", "Reconhecimento em tratamento de feridas", "03", PURPLE),
        ("Agenda Preenchida", "Faturamento mais previsível e estável", "04", AMBER),
        ("ROI Positivo", "Retorno estimado entre 60 e 90 dias", "05", GREEN),
    ]

    for i, (title, desc, num, accent) in enumerate(results):
        y = Inches(1.4) + i * Inches(0.95)
        _add_round_rect(s, Inches(0.7), y, Inches(11.9), Inches(0.85), BG_CARD, 0.08)
        _add_rect(s, Inches(0.7), y, Inches(0.1), Inches(0.85), accent)
        _textbox(s, Inches(1.1), y + Inches(0.2), Inches(0.7), Inches(0.45), num, size=18, color=accent, bold=True)
        _textbox(s, Inches(2.0), y + Inches(0.12), Inches(9.5), Inches(0.35), title, size=16, color=WHITE, bold=True)
        _textbox(s, Inches(2.0), y + Inches(0.45), Inches(9.5), Inches(0.3), desc, size=12, color=SLATE_400)


def slide_por_que_code(prs):
    s = _new_slide(prs)
    _section_title(s, "Por que a CODE?",
                   "Tecnologia própria, foco em conversão e compromisso com resultado")
    _footer(s, 11)

    reasons = [
        ("+300 empresas transformadas", "Track record comprovado em múltiplos segmentos"),
        ("Plataformas próprias", "Nexus CRM, ERP e soluções com Inteligência Artificial"),
        ("Foco em métricas e ROI", "Decisões baseadas em dados, conversão e retorno"),
        ("Garantia de qualidade", "Suporte contínuo e prazos cumpridos"),
    ]

    for i, (title, desc) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * Inches(6.2)
        y = Inches(1.45) + row * Inches(1.7)
        _add_round_rect(s, x, y, Inches(5.9), Inches(1.45), BG_CARD, 0.06)
        _add_rect(s, x, y, Inches(0.1), Inches(1.45), TEAL)
        _textbox(s, x + Inches(0.4), y + Inches(0.3), Inches(5.2), Inches(0.4),
                 f"✓  {title}", size=16, color=WHITE, bold=True)
        _textbox(s, x + Inches(0.7), y + Inches(0.8), Inches(4.9), Inches(0.4),
                 desc, size=13, color=SLATE_400)

    # Quote
    _add_round_rect(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.5), BG_CARD_ALT, 0.05)
    _textbox(s, Inches(1.1), Inches(5.4), Inches(11), Inches(0.9),
             '"A CODE não entrega apenas um site.\nEntregamos uma máquina de geração de pacientes."',
             size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)


def slide_cta(prs):
    s = _new_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.08), TEAL)
    _add_rect(s, 0, 0, Inches(0.12), SLIDE_H, TEAL)

    _pill(s, Inches(4.9), Inches(1.3), Inches(3.5), Inches(0.4),
          "VAMOS COMEÇAR?", bg=TEAL, fg=BG, size=12)

    _textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.7),
             "CODE Tecnologia Empresarial", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(1.5), Inches(2.75), Inches(10.3), Inches(0.6),
             "Soluções completas para automatizar, vender mais e crescer com inteligência.",
             size=15, color=SLATE_400, align=PP_ALIGN.CENTER)

    contacts = [
        ("Site", "codetechoficial.com.br"),
        ("WhatsApp", "(14) 99690-2902"),
        ("E-mail", "contato@codetechoficial.com.br"),
    ]

    for i, (label, value) in enumerate(contacts):
        x = Inches(1.3) + i * Inches(3.9)
        _add_round_rect(s, x, Inches(3.8), Inches(3.5), Inches(1.5), BG_CARD, 0.08)
        _textbox(s, x + Inches(0.2), Inches(4.05), Inches(3.1), Inches(0.35),
                 label, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.2), Inches(4.5), Inches(3.1), Inches(0.5),
                 value, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _textbox(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.5),
             "Dra. Jussara Bertoloto  ·  Projeto de Presença Digital",
             size=13, color=SLATE_500, align=PP_ALIGN.CENTER)
    _textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4),
             "Landing Page Premium  ·  Tráfego Pago  ·  SEO Local  ·  Infraestrutura CODE",
             size=12, color=SLATE_500, align=PP_ALIGN.CENTER)



def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Capa
    slide_cover(prs)

    # 2 — Roteiro
    s = _new_slide(prs)
    _section_title(s, "Roteiro da Apresentação", "Visão geral do conteúdo desta proposta")
    _footer(s, 2)
    items = [
        ("01", "Quem Somos — CODE"),
        ("02", "O Problema de Mercado"),
        ("03", "Produtos & Serviços da Dra. Jussara"),
        ("04", "A Solução CODE (LP + Estratégia)"),
        ("05", "Estrutura da Landing Page"),
        ("06", "Estratégia de Tráfego & Visibilidade"),
        ("07", "Infraestrutura Técnica"),
        ("08", "Resultados Esperados (90 dias)"),
        ("09", "Por que a CODE"),
        ("10", "Vamos Começar — Contato"),
    ]
    for i, (num, title) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.9) + col * Inches(6.0)
        y = Inches(1.45) + row * Inches(0.9)
        _textbox(s, x, y, Inches(0.7), Inches(0.45), num, size=18, color=TEAL, bold=True)
        _textbox(s, x + Inches(0.85), y + Inches(0.05), Inches(4.8), Inches(0.4),
                 title, size=15, color=WHITE)

    # 3–12
    slide_quem_somos(prs)
    slide_problema(prs)
    slide_servicos(prs)
    slide_solucao(prs)
    slide_estrutura_lp(prs)
    slide_trafego(prs)
    slide_infra(prs)
    slide_resultados(prs)
    slide_por_que_code(prs)
    slide_cta(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    ARTIFACT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ARTIFACT))
    print(f"Saved: {OUT_FILE}")
    print(f"Artifact: {ARTIFACT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
